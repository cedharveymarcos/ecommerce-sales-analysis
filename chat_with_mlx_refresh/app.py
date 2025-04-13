import argparse
import atexit
import base64
import copy
import hashlib
import logging
import re
import threading
import time
from pathlib import Path
from typing import Dict, List, Tuple, Optional, Union

import chromadb
import gradio as gr
from gradio.components.chatbot import ChatMessage
from langchain_text_splitters import RecursiveCharacterTextSplitter
from pandas import DataFrame
from sentence_transformers import SentenceTransformer

from .language import get_text
from .model import Message, MessageRole, ModelManager, TextModel, VisionModel, OpenAIModel

logger = logging.getLogger(__name__)
logger.setLevel(logging.INFO)

model_manager = ModelManager()
generation_stop_event = threading.Event()


def get_loaded_model() -> Union[TextModel, VisionModel, OpenAIModel]:
    model = model_manager.get_loaded_model()
    if model is None:
        raise RuntimeError("No model loaded.")
    return model


def get_file_md5(file_name: Path) -> str:
    md5 = hashlib.md5()
    with open(file_name, "rb") as f:
        for chunk in iter(lambda: f.read(4096), b""):
            md5.update(chunk)
    return md5.hexdigest()


class FileManager:
    def __init__(self):
        self.files = {}

    def format_content(self, file_name, content):
        boundary_start = f"<<<BEGIN FILE:{file_name}>>>"
        boundary_end = f"<<<END FILE:{file_name}>>>"
        content = content.replace(boundary_start, '')
        content = content.replace(boundary_end, '')
        return f"{boundary_start}\n{content}\n{boundary_end}"

    def load_file(self, file_name: Path, raw_content_only: bool = False):
        if file_name.name in self.files:
            if self.files[file_name.name]["md5"] == get_file_md5(file_name):
                return self.files[file_name.name]["formatted_content"] if not raw_content_only else self.files[file_name.name]["content"]
        suffix = file_name.suffix.lower()
        if suffix == ".pdf":
            return self.load_pdf(file_name, raw_content_only)
        elif suffix in [".txt", ".csv", ".md"]:
            return self.load_txt_like(file_name, raw_content_only)
        elif suffix == ".docx":
            return self.load_docx(file_name, raw_content_only)
        elif suffix == ".pptx":
            return self.load_pptx(file_name, raw_content_only)
        elif suffix in [".xlsx", ".xls"]:
            return self.load_excel(file_name, raw_content_only)
        else:
            return None

    def load_pdf(self, file_name: Path, raw_content_only: bool = False):
        try:
            from pypdf import PdfReader

            pdf = PdfReader(file_name)
            content_parts = []
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    content_parts.append(text)
            content = ''.join(content_parts)
            formatted_content = self.format_content(file_name, content)
            self.files[file_name.name] = {
                "md5": get_file_md5(file_name),
                "formatted_content": formatted_content,
                "content": content
            }
            return formatted_content if not raw_content_only else content
        except ImportError:
            logger.warning("pypdf not found.")
            return None

    def load_txt_like(self, file_name: Path, raw_content_only: bool = False):
        with open(file_name, 'r', encoding='utf-8') as f:
            content = f.read()
        formatted_content = self.format_content(file_name, content)
        self.files[file_name.name] = {
            "md5": get_file_md5(file_name),
            "formatted_content": formatted_content,
            "content": content
        }
        return formatted_content if not raw_content_only else content

    def load_docx(self, file_name: Path, raw_content_only: bool = False):
        try:
            from docx import Document

            doc = Document(str(file_name))
            content = "\n".join([para.text for para in doc.paragraphs])
            formatted_content = self.format_content(file_name, content)
            self.files[file_name.name] = {
                "md5": get_file_md5(file_name),
                "formatted_content": formatted_content,
                "content": content
            }
            return formatted_content if not raw_content_only else content
        except ImportError:
            logger.warning("docx not found.")
            return None

    def load_pptx(self, file_name: Path, raw_content_only: bool = False):
        try:
            from pptx import Presentation

            prs = Presentation(str(file_name))
            content_parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        content_parts.append(shape.text)
            content = "\n".join(content_parts)
            formatted_content = self.format_content(file_name, content)
            self.files[file_name.name] = {
                "md5": get_file_md5(file_name),
                "formatted_content": formatted_content,
                "content": content
            }
            return formatted_content if not raw_content_only else content
        except ImportError:
            logger.warning("pptx not found.")
            return None

    def load_excel(self, file_name: Path, raw_content_only: bool = False):
        try:
            from pandas import ExcelFile, read_excel
            excel_file = ExcelFile(file_name)
            content_parts = []

            for sheet_name in excel_file.sheet_names:
                df = read_excel(excel_file, sheet_name=sheet_name)
                content_parts.append(f"Sheet: {sheet_name}\n")
                content_parts.append(df.to_csv(index=False))

            content = ''.join(content_parts)
            formatted_content = self.format_content(file_name, content)
            self.files[file_name.name] = {
                "md5": get_file_md5(file_name),
                "formatted_content": formatted_content,
                "content": content
            }
            return formatted_content if not raw_content_only else content
        except ImportError:
            logger.warning("pandas not found.")
            return None

    def clear(self):
        self.files.clear()


class RAGManager:
    def __init__(self, model_name: str = "all-MiniLM-L6-v2", persistence=False, db_path="./chromadb"):
        if persistence:
            self.client = chromadb.PersistentClient(db_path, settings=chromadb.Settings(anonymized_telemetry=False))
        else:
            self.client = chromadb.Client(settings=chromadb.Settings(anonymized_telemetry=False))

        self.embedding_model = SentenceTransformer(model_name)
        self.chunk_size = 1000
        self.chunk_overlap = 200
        self.n_results = 5
        self.similarity_threshold = 0.0

        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=self.chunk_size,
            chunk_overlap=self.chunk_overlap
        )
        self.collection = self.client.get_or_create_collection(
            name="rag_collection",
            metadata={"hnsw:space": "cosine"}
        )
        self.indexed_files = set()
        self.enabled = False

    def update_parameters(self,
                          chunk_size: int = None,
                          chunk_overlap: int = None,
                          n_results: int = None,
                          similarity_threshold: float = None):
        updated = False

        if chunk_size is not None and chunk_size != self.chunk_size:
            self.chunk_size = chunk_size
            updated = True

        if chunk_overlap is not None and chunk_overlap != self.chunk_overlap:
            self.chunk_overlap = chunk_overlap
            updated = True

        if n_results is not None:
            self.n_results = n_results

        if similarity_threshold is not None:
            self.similarity_threshold = similarity_threshold

        if updated:
            self.text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=self.chunk_size,
                chunk_overlap=self.chunk_overlap
            )

        return updated

    def add_document(self, file_path: Path, file_content: str) -> Tuple[bool, str]:
        if file_path.name in self.indexed_files:
            return False, "Document '{}' is already indexed.".format(file_path.name)

        chunks = self.text_splitter.split_text(file_content)
        if not chunks:
            return False, "No content to index in document '{}'.".format(file_path.name)

        embeddings = self.embedding_model.encode(chunks, show_progress_bar=True)

        ids = [f"{file_path.name}_{i}" for i in range(len(chunks))]
        metadata = [{"source": file_path.name, "chunk_id": i} for i in range(len(chunks))]

        self.collection.add(
            embeddings=embeddings,
            documents=chunks,
            metadatas=metadata,
            ids=ids
        )
        self.indexed_files.add(file_path.name)
        return True, "Document '{}' indexed with {} chunks.".format(file_path.name, len(chunks))

    def retrieve(self, query: str, n_results: int = None) -> str:
        if self.collection.count() == 0:
            return ""

        results_count = n_results if n_results is not None else self.n_results

        query_embeddings = self.embedding_model.encode([query])
        results = self.collection.query(
            query_embeddings=query_embeddings,
            n_results=min(results_count, self.collection.count())
        )

        if not results or not results['documents'] or not results['documents'][0]:
            return ""

        filtered_docs = []
        documents = results['documents'][0]
        distances = results.get('distances', [None])[0] if 'distances' in results else [None] * len(documents)

        for doc, distance in zip(documents, distances):
            if distance is not None and distance >= self.similarity_threshold:
                filtered_docs.append(doc)

        if not filtered_docs:
            return ""

        retrieved_docs = "\n\n".join(filtered_docs)
        return retrieved_docs

    def clear_index(self):
        self.client.delete_collection(name=self.collection.name)
        self.collection = self.client.get_or_create_collection(
            name="rag_collection",
            metadata={"hnsw:space": "cosine"}
        )
        self.indexed_files.clear()

    def enable(self):
        self.enabled = True

    def disable(self):
        self.enabled = False

    def is_enabled(self) -> bool:
        return self.enabled

    def get_status(self) -> Tuple[bool, str]:
        if not self.indexed_files:
            return False, "RAG Index is empty."
        total_chunks = self.collection.count()
        return True, "Indexed {} documents with {} chunks.".format(len(self.indexed_files), total_chunks)

    def get_parameters(self) -> Dict:
        return {
            "chunk_size": self.chunk_size,
            "chunk_overlap": self.chunk_overlap,
            "n_results": self.n_results,
            "similarity_threshold": self.similarity_threshold
        }


file_manager = FileManager()
rag_manager = RAGManager()


def preprocess_file(message: Dict, history: List[Dict]) -> Tuple[str, List[Dict]]:
    processed_message_parts = []
    if "files" in message and message["files"]:
        for file_path_str in message["files"]:
            if file_path_str:
                file_content = file_manager.load_file(Path(file_path_str))
                if file_content:
                    processed_message_parts.append(file_content)

    text_content = message.get("text", "")
    if not isinstance(text_content, str):
        text_content = str(text_content) if text_content is not None else ""
    processed_message_parts.append(text_content)

    processed_message = "".join(processed_message_parts)

    preprocessed_history = []
    i = 0
    while i < len(history):
        current_hist_item_original = history[i]
        current_processed_item = copy.deepcopy(current_hist_item_original)

        current_content = current_hist_item_original.get("content")

        if isinstance(current_content, tuple):
            combined_files_text_parts = []
            for file_path_in_tuple_str in current_content:
                if file_path_in_tuple_str:
                    file_actual_content = file_manager.load_file(Path(file_path_in_tuple_str))
                    if file_actual_content:
                        combined_files_text_parts.append(file_actual_content)

            final_combined_content = "".join(combined_files_text_parts)

            text_to_merge_from_next = ""
            consumed_next_item_flag = False
            if i + 1 < len(history):
                next_original_hist_item = history[i + 1]
                next_content = next_original_hist_item.get("content")
                if isinstance(next_content, str):
                    text_to_merge_from_next = next_content
                    consumed_next_item_flag = True

            current_processed_item["content"] = final_combined_content + text_to_merge_from_next
            preprocessed_history.append(current_processed_item)

            i += 1
            if consumed_next_item_flag:
                i += 1
        else:
            if not isinstance(current_content, str):
                current_processed_item["content"] = str(current_content) if current_content is not None else ""

            preprocessed_history.append(current_processed_item)
            i += 1

    return processed_message, preprocessed_history


def encode_image(image_path: Path):
    suffix = image_path.suffix.lower()
    mime_types = {
        ".png": "image/png",
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
        ".heif": "image/heif",
        ".heic": "image/heic"
    }
    if suffix not in mime_types:
        raise RuntimeError(f"Unsupported imgae type: {suffix}")
    with open(image_path, 'rb') as f:
        return f"data:{mime_types[suffix]};base64,{base64.b64encode(f.read()).decode('utf-8')}"


def prepare_openai_message_content(text_input: Optional[str], file_paths: Optional[List[str]]) -> List[Dict]:
    content_parts = []
    current_text = str(text_input) if text_input is not None else ""

    if file_paths:
        for file_path_str in file_paths:
            file = Path(file_path_str)
            if not file.is_file():
                logger.warning(f"File not found: {file_path_str}, skipping.")
                continue

            suffix = file.suffix.lower()
            if suffix in [".png", ".jpg", ".jpeg", ".heif", ".heic"]:
                try:
                    image_base64 = encode_image(file)
                    content_parts.append({"type": "image_url", "image_url": {"url": image_base64}})
                except Exception as e:
                    logger.error(f"Failed to encode image {file_path_str}: {e}")
            else:
                file_content = file_manager.load_file(file)
                if file_content:
                    current_text += "\n" + file_content
                else:
                    logger.warning(f"Could not load content from file: {file_path_str}")

    if current_text or not content_parts:
        content_parts.insert(0, {"type": "text", "text": current_text.strip()})

    return content_parts


def build_openai_api_messages(current_message_dict: Dict, history_list: List[Dict], system_prompt: Optional[str], ) -> List[Dict]:
    api_messages = []
    if system_prompt and system_prompt.strip():
        api_messages.append({"role": "system", "content": system_prompt})

    i = 0
    while i < len(history_list):
        hist_item = history_list[i]
        role = hist_item.get("role")
        content = hist_item.get("content")

        if role not in ("user", "assistant", "tool"):
            logger.warning(f"Skipping history item with unexpected role: {role}")
            i += 1
            continue

        if isinstance(content, tuple):
            files_in_hist_item = list(content)
            text_for_files = None
            if i + 1 < len(history_list) and \
                    history_list[i + 1].get("role") == role and \
                    isinstance(history_list[i + 1].get("content"), str):
                text_for_files = history_list[i + 1]["content"]
                i += 1

            processed_content_parts = prepare_openai_message_content(text_for_files, files_in_hist_item)
        elif isinstance(content, str):
            processed_content_parts = prepare_openai_message_content(content, None)
        elif isinstance(content, list):
            processed_content_parts = content
        else:
            logger.warning(f"Skipping history item with unexpected content type: {type(content)}")
            i += 1
            continue

        if processed_content_parts:
            api_messages.append({"role": role, "content": processed_content_parts})
        i += 1

    user_text = current_message_dict.get("text")
    user_files = current_message_dict.get("files")
    current_user_content_parts = prepare_openai_message_content(user_text, user_files)
    if current_user_content_parts:
        api_messages.append({"role": "user", "content": current_user_content_parts})

    return api_messages


def prepare_generic_model_inputs(current_message_dict: Dict, history_list: List[Dict], system_prompt: Optional[str], model_instance: Union[TextModel, VisionModel]) -> Tuple[str, List[Dict], List[str]]:
    effective_history = []
    if system_prompt and system_prompt.strip() != "":
        effective_history.append(Message(MessageRole.SYSTEM, content=system_prompt).to_dict())
    effective_history.extend(history_list)

    effective_history = [
        history_item for history_item in effective_history
        if not (isinstance(history_item, dict) and
                isinstance(history_item.get("metadata"), dict) and
                history_item["metadata"].get("title") in ["Thought for", "Thinking"])
    ]

    image_paths = []
    if isinstance(model_instance, VisionModel):
        for hist_item in effective_history:
            files_to_check = []
            if isinstance(hist_item.get("content"), tuple):
                files_to_check.extend(list(hist_item["content"]))
            elif isinstance(hist_item.get("files"), list):
                files_to_check.extend(hist_item["files"])

            for file_path_str in files_to_check:
                path = Path(file_path_str)
                if path.is_file() and path.suffix.lower() in [".jpg", ".png", ".jpeg"]:
                    image_paths.append(file_path_str)

        if "files" in current_message_dict and current_message_dict["files"]:
            for file_path_str in current_message_dict["files"]:
                path = Path(file_path_str)
                if path.is_file() and path.suffix.lower() in [".jpg", ".png", ".jpeg"]:
                    image_paths.append(file_path_str)

        image_paths = list(set(image_paths))

    processed_message_text, processed_history_list = preprocess_file(current_message_dict, effective_history)

    return processed_message_text, processed_history_list, image_paths


def handle_chat(message: Dict,
                history: List[Dict],
                system_prompt: str = None,
                temperature: float = 0.7,
                top_k: int = 20,
                top_p: float = 0.9,
                min_p: float = 0.0,
                max_tokens: int = 512,
                repetition_penalty: float = 1.0,
                rag_enabled: bool = False,
                rag_n_results: int = 5,
                stream: bool = True):
    if rag_enabled and message.get("text"):
        original_text = message.get("text", "")
        enhanced_text = enhance_message_with_rag(original_text, rag_enabled, rag_n_results)
        message = dict(message)
        message["text"] = enhanced_text

    CHATML_CONTROL_TOKENS = [
        '<|im_start|>', '<|im_end|>',
        '<|system|>', '<|user|>', '<|assistant|>',
        '<|end|>', '<|endoftext|>'
    ]

    CHATML_STOP_TOKENS = ['<|im_end|>', '<|end|>', '<|endoftext|>']

    def filter_chatml_tokens(text: str) -> tuple[str, bool]:
        should_stop = False
        filtered_text = text

        for stop_token in CHATML_STOP_TOKENS:
            if stop_token in text:
                should_stop = True
                stop_index = text.find(stop_token)
                filtered_text = text[:stop_index]
                break

        for token in CHATML_CONTROL_TOKENS:
            filtered_text = filtered_text.replace(token, '')

        return filtered_text, should_stop

    try:
        model = get_loaded_model()

        if isinstance(model, OpenAIModel):
            api_messages = build_openai_api_messages(message, history, system_prompt)

            response_stream = model.generate_response(messages=api_messages, stream=stream, temperature=temperature, top_p=top_p, max_tokens=max_tokens, n=1)

            chat_message_accumulator = ChatMessage(role="assistant", content="")
            thinking_message = None
            thinking_content_parts = []
            final_content_parts = []
            in_thinking = False
            thinking_start_time = None
            chunk_buffer = ""

            if stream:
                for chunk in response_stream:
                    if generation_stop_event.is_set():
                        break
                    if chunk.choices:
                        delta = chunk.choices[0].delta
                        if delta.content:
                            chunk_text = ''.join([chunk_buffer, delta.content])
                            chunk_buffer = ""

                            for partial in ["<", "<t", "<th", "<thi", "<thin", "<think", "</", "</t", "</th", "</thi", "</thin", "</think"]:
                                if chunk_text.endswith(partial):
                                    chunk_buffer = partial
                                    chunk_text = chunk_text[:-len(partial)]
                                    break

                            if not chunk_text:
                                continue

                            chunk_text, should_stop = filter_chatml_tokens(chunk_text)

                            if should_stop:
                                if in_thinking:
                                    thinking_content_parts.append(chunk_text)
                                    thinking_message.content = ''.join(thinking_content_parts)
                                    thinking_message.metadata["title"] = "Thought for"
                                    thinking_message.metadata["status"] = "done"
                                    thinking_message.metadata["duration"] = time.time() - thinking_start_time
                                    yield [thinking_message, chat_message_accumulator]
                                else:
                                    final_content_parts.append(chunk_text)
                                    chat_message_accumulator.content = ''.join(final_content_parts)
                                    if thinking_message and thinking_message.metadata.get("status") == "done":
                                        yield [thinking_message, chat_message_accumulator]
                                    else:
                                        yield chat_message_accumulator
                                break

                            if "<think>" in chunk_text and not in_thinking:
                                in_thinking = True
                                thinking_start_time = time.time()
                                thinking_message = ChatMessage(
                                    role="assistant",
                                    content="",
                                    metadata={"title": "Thinking", "id": 0, "status": "pending"}
                                )
                                think_start_idx = chunk_text.find("<think>") + len("<think>")
                                if think_start_idx < len(chunk_text):
                                    thinking_content_parts.append(chunk_text[think_start_idx:])
                                    thinking_message.content = ''.join(thinking_content_parts)
                                yield thinking_message
                                continue

                            elif in_thinking and "</think>" not in chunk_text:
                                thinking_content_parts.append(chunk_text)
                                thinking_message.content = ''.join(thinking_content_parts)
                                yield thinking_message
                                continue

                            elif in_thinking and "</think>" in chunk_text:
                                think_end_idx = chunk_text.find("</think>")
                                thinking_content_parts.append(chunk_text[:think_end_idx])
                                thinking_message.content = ''.join(thinking_content_parts)
                                thinking_message.metadata["title"] = "Thought for"
                                thinking_message.metadata["status"] = "done"
                                thinking_message.metadata["duration"] = time.time() - thinking_start_time

                                remaining_content = chunk_text[think_end_idx + len("</think>"):]
                                if remaining_content.strip():
                                    final_content_parts.append(remaining_content)
                                    chat_message_accumulator.content = ''.join(final_content_parts)

                                yield [thinking_message, chat_message_accumulator]
                                in_thinking = False
                                continue

                            elif not in_thinking:
                                if "<think>" in chunk_text and "</think>" in chunk_text:
                                    think_match = re.search(r'<think>(.*?)</think>', chunk_text, re.DOTALL)
                                    if think_match:
                                        thinking_start_time = time.time()
                                        extracted_thinking = think_match.group(1)
                                        thinking_message = ChatMessage(
                                            role="assistant",
                                            content=extracted_thinking,
                                            metadata={
                                                "title": "Thought for",
                                                "id": 0,
                                                "status": "done",
                                                "duration": 0.1
                                            }
                                        )

                                        before_think = chunk_text[:chunk_text.find("<think>")]
                                        after_think = chunk_text[chunk_text.find("</think>") + len("</think>"):]
                                        final_content_parts.extend([before_think, after_think])
                                        chat_message_accumulator.content = ''.join(final_content_parts)

                                        yield [thinking_message, chat_message_accumulator]
                                        continue
                                else:
                                    final_content_parts.append(chunk_text)
                                    chat_message_accumulator.content = ''.join(final_content_parts)
                                    if thinking_message and thinking_message.metadata.get("status") == "done":
                                        yield [thinking_message, chat_message_accumulator]
                                    else:
                                        yield chat_message_accumulator
            else:
                if response_stream.choices:
                    full_content = response_stream.choices[0].message.content

                    full_content, _ = filter_chatml_tokens(full_content)

                    think_match = re.search(r'<think>(.*?)</think>', full_content, re.DOTALL)
                    if think_match:
                        thinking_content = think_match.group(1)
                        thinking_message = ChatMessage(
                            role="assistant",
                            content=thinking_content,
                            metadata={
                                "title": "Thought for",
                                "id": 0,
                                "status": "done",
                                "duration": 0.1
                            }
                        )

                        final_content = re.sub(r'<think>.*?</think>', '', full_content, flags=re.DOTALL).strip()
                        chat_message_accumulator.content = final_content

                        yield [thinking_message, chat_message_accumulator]
                    else:
                        chat_message_accumulator.content = full_content
                        yield chat_message_accumulator
                else:
                    logger.error("OpenAI non-stream response had no choices.")
                    yield ChatMessage(role="assistant", content="Error: No response from model.")
        else:
            tokenizer_to_check = None
            if isinstance(model, VisionModel):
                if model.processor and hasattr(model.processor, 'tokenizer'):
                    tokenizer_to_check = model.processor.tokenizer
            elif hasattr(model, 'tokenizer'):
                tokenizer_to_check = model.tokenizer

            if tokenizer_to_check and tokenizer_to_check.chat_template is None:
                raise RuntimeError("Model {} does not have a chat template. Please use the 'Completion' tab or set a chat template.".format(model.model_name if hasattr(model, 'model_name') else type(model).__name__))

            processed_message_text, processed_history_list, image_paths = prepare_generic_model_inputs(message, history, system_prompt, model)

            temperature = float(temperature)
            top_k = int(top_k)
            top_p = float(top_p)
            min_p = float(min_p)
            repetition_penalty = float(repetition_penalty)

            response_args = {
                "message": processed_message_text,
                "history": processed_history_list,
                "stream": stream,
                "temperature": temperature,
                "top_k": top_k,
                "top_p": top_p,
                "min_p": min_p,
                "max_tokens": max_tokens,
                "repetition_penalty": repetition_penalty,
            }
            if isinstance(model, VisionModel):
                response_args["images"] = image_paths
                eos_token = model.processor.tokenizer.eos_token if model.processor and model.processor.tokenizer else None
            else:
                eos_token = model.tokenizer.eos_token if model.tokenizer else None

            response_stream = model.generate_response(**response_args)

            chat_message_accumulator = ChatMessage(role="assistant", content="")
            thinking_message = None
            thinking_content_parts = []
            final_content_parts = []
            in_thinking = False
            thinking_start_time = None
            chunk_buffer = ""

            for chunk in response_stream:
                if generation_stop_event.is_set():
                    break
                chunk_text = ""
                if isinstance(chunk, str):
                    chunk_text = chunk
                elif hasattr(chunk, "text"):
                    chunk_text = chunk.text
                elif hasattr(chunk, "choices") and chunk.choices and hasattr(chunk.choices[0], "delta") and chunk.choices[0].delta.content:
                    chunk_text = chunk.choices[0].delta.content
                else:
                    logger.warning(f"Unexpected chunk type from model: {type(chunk)}")

                if chunk_text:
                    if stream and eos_token and eos_token in chunk_text:
                        if chunk_text == eos_token:
                            break
                        chunk_text = chunk_text.split(eos_token)[0]

                    chunk_text = ''.join([chunk_buffer, chunk_text])
                    chunk_buffer = ""

                    for partial in ["<", "<t", "<th", "<thi", "<thin", "<think", "</", "</t", "</th", "</thi", "</thin", "</think"]:
                        if chunk_text.endswith(partial):
                            chunk_buffer = partial
                            chunk_text = chunk_text[:-len(partial)]
                            break

                    if not chunk_text:
                        continue

                    chunk_text, should_stop = filter_chatml_tokens(chunk_text)

                    if should_stop:
                        if in_thinking:
                            thinking_content_parts.append(chunk_text)
                            thinking_message.content = ''.join(thinking_content_parts)
                            thinking_message.metadata["title"] = "Thought for"
                            thinking_message.metadata["status"] = "done"
                            thinking_message.metadata["duration"] = time.time() - thinking_start_time
                            if stream:
                                yield [thinking_message, chat_message_accumulator]
                        else:
                            final_content_parts.append(chunk_text)
                            chat_message_accumulator.content = ''.join(final_content_parts)
                            if thinking_message and thinking_message.metadata.get("status") == "done":
                                if stream:
                                    yield [thinking_message, chat_message_accumulator]
                            else:
                                if stream:
                                    yield chat_message_accumulator
                        break

                    if "<think>" in chunk_text and not in_thinking:
                        in_thinking = True
                        thinking_start_time = time.time()
                        thinking_message = ChatMessage(
                            role="assistant",
                            content="",
                            metadata={"title": "Thinking", "id": 0, "status": "pending"}
                        )
                        think_start_idx = chunk_text.find("<think>") + len("<think>")
                        if think_start_idx < len(chunk_text):
                            thinking_content_parts.append(chunk_text[think_start_idx:])
                            thinking_message.content = ''.join(thinking_content_parts)
                        if stream:
                            yield thinking_message
                        continue

                    elif in_thinking and "</think>" not in chunk_text:
                        thinking_content_parts.append(chunk_text)
                        thinking_message.content = ''.join(thinking_content_parts)
                        if stream:
                            yield thinking_message
                        continue

                    elif in_thinking and "</think>" in chunk_text:
                        think_end_idx = chunk_text.find("</think>")
                        thinking_content_parts.append(chunk_text[:think_end_idx])
                        thinking_message.content = ''.join(thinking_content_parts)
                        thinking_message.metadata['title'] = "Thought for"
                        thinking_message.metadata["status"] = "done"
                        thinking_message.metadata["duration"] = time.time() - thinking_start_time

                        remaining_content = chunk_text[think_end_idx + len("</think>"):]
                        if remaining_content.strip():
                            final_content_parts.append(remaining_content)
                            chat_message_accumulator.content = ''.join(final_content_parts)

                        if stream:
                            yield [thinking_message, chat_message_accumulator]
                        in_thinking = False
                        continue

                    elif not in_thinking:
                        if "<think>" in chunk_text and "</think>" in chunk_text:
                            think_match = re.search(r'<think>(.*?)</think>', chunk_text, re.DOTALL)
                            if think_match:
                                thinking_start_time = time.time()
                                extracted_thinking = think_match.group(1)
                                thinking_message = ChatMessage(
                                    role="assistant",
                                    content=extracted_thinking,
                                    metadata={
                                        "title": "Thought for",
                                        "id": 0,
                                        "status": "done",
                                        "duration": 0.1
                                    }
                                )

                                before_think = chunk_text[:chunk_text.find("<think>")]
                                after_think = chunk_text[chunk_text.find("</think>") + len("</think>"):]
                                final_content_parts.extend([before_think, after_think])
                                chat_message_accumulator.content = ''.join(final_content_parts)

                                if stream:
                                    yield [thinking_message, chat_message_accumulator]
                                continue
                        else:
                            final_content_parts.append(chunk_text)
                            chat_message_accumulator.content = ''.join(final_content_parts)
                            if thinking_message and thinking_message.metadata.get("status") == "done":
                                if stream:
                                    yield [thinking_message, chat_message_accumulator]
                            else:
                                if stream:
                                    yield chat_message_accumulator

                    if stream and eos_token and eos_token in chunk_text:
                        break

            if not stream:
                if thinking_message and thinking_message.metadata.get("status") == "done":
                    yield [thinking_message, chat_message_accumulator]
                else:
                    yield chat_message_accumulator
    except Exception as e:
        logger.exception("Error in handle_chat:")
        raise gr.Error(str(e))


def managed_chat_generator(
        message: Dict,
        history: List[Dict],
        system_prompt: str = None,
        temperature: float = 0.7,
        top_k: int = 20,
        top_p: float = 0.9,
        min_p: float = 0.0,
        max_tokens: int = 512,
        repetition_penalty: float = 1.0,
        rag_enabled: bool = False,
        rag_n_results: int = 5,
        stream: bool = True):
    g = handle_chat(
        message=message,
        history=history,
        system_prompt=system_prompt,
        temperature=temperature,
        top_k=top_k,
        top_p=top_p,
        min_p=min_p,
        max_tokens=max_tokens,
        repetition_penalty=repetition_penalty,
        rag_enabled=rag_enabled,
        rag_n_results=rag_n_results,
        stream=stream,
    )

    generation_stop_event.clear()
    model_manager.close_active_generator()
    model_manager.set_active_generator(g)

    try:
        yield from g
    finally:
        model_manager.close_active_generator()


def handle_completion(prompt: str,
                      temperature: float = 0.7,
                      top_k: int = 20,
                      top_p: float = 0.9,
                      min_p: float = 0.0,
                      max_tokens: int = 512,
                      repetition_penalty: float = 1.0,
                      stream: bool = True):
    try:
        model = get_loaded_model()
        if isinstance(model, VisionModel) or isinstance(model, OpenAIModel):
            raise RuntimeError("Not supported yet.")

        temperature = float(temperature)
        top_p = float(top_p)
        repetition_penalty = float(repetition_penalty)

        if not stream:
            completion_text = model.generate_completion(
                prompt=prompt,
                stream=stream,
                temperature=temperature,
                top_k=top_k,
                top_p=top_p,
                min_p=min_p,
                max_tokens=max_tokens,
                repetition_penalty=repetition_penalty)
            return ''.join([prompt, completion_text])
        else:
            response_parts = [prompt]
            eos_token = model.tokenizer.eos_token

            for chunk in model.generate_completion(
                    prompt=prompt,
                    stream=stream,
                    temperature=temperature,
                    top_k=top_k,
                    top_p=top_p,
                    min_p=min_p,
                    max_tokens=max_tokens,
                    repetition_penalty=repetition_penalty):
                if generation_stop_event.is_set():
                    break

                chunk_text = ""
                if isinstance(chunk, str):
                    chunk_text = chunk
                elif hasattr(chunk, "text"):
                    chunk_text = chunk.text
                elif hasattr(chunk, "choices") and chunk.choices and hasattr(chunk.choices[0], "delta") and chunk.choices[0].delta.content:
                    chunk_text = chunk.choices[0].delta.content
                else:
                    logger.warning(f"Unexpected chunk type from model: {type(chunk)}")

                if chunk_text:
                    if eos_token not in chunk_text:
                        response_parts.append(chunk_text)
                        yield ''.join(response_parts)
                    else:
                        if eos_token in chunk_text:
                            before_eos = chunk_text.split(eos_token)[0]
                            if before_eos:
                                response_parts.append(before_eos)
                        yield ''.join(response_parts)
                        break
            return None
    except Exception as e:
        raise gr.Error(str(e))


def managed_completion_generator(prompt: str,
                                 temperature: float = 0.7,
                                 top_k: int = 20,
                                 top_p: float = 0.9,
                                 min_p: float = 0.0,
                                 max_tokens: int = 512,
                                 repetition_penalty: float = 1.0,
                                 stream: bool = True):
    g = handle_completion(
        prompt=prompt,
        temperature=temperature,
        top_k=top_k,
        top_p=top_p,
        min_p=min_p,
        max_tokens=max_tokens,
        repetition_penalty=repetition_penalty,
        stream=stream
    )

    generation_stop_event.clear()
    model_manager.close_active_generator()
    model_manager.set_active_generator(g)

    try:
        yield from g
    finally:
        model_manager.close_active_generator()


def get_load_model_status():
    if model_manager.get_loaded_model_config():
        return get_text("Page.Chat.LoadModelBlock.Textbox.model_status.loaded_value").format(model_manager.get_loaded_model_config().get("display_name"))
    else:
        return get_text("Page.Chat.LoadModelBlock.Textbox.model_status.not_loaded_value")


def load_model(model_name: str) -> Tuple[str, str]:
    try:
        model_manager.load_model(model_name)
        return get_load_model_status(), model_manager.get_system_prompt(default=True)
    except Exception as e:
        raise gr.Error(str(e))


def stop_generation() -> None:
    generation_stop_event.set()


def chat_load_model_callback(model_name: str):
    stop_generation()
    return load_model(model_name)


def completion_load_model_callback(model_name: str):
    stop_generation()
    return load_model(model_name)[0]


def get_default_system_prompt_callback():
    if model_manager.get_loaded_model_config():
        return model_manager.get_system_prompt(default=True)
    else:
        raise gr.Error("No model loaded.")


def update_model_management_models_list():
    return DataFrame({get_text("Page.ModelManagement.Dataframe.model_list.headers"): model_manager.get_model_list()})


def update_select_model_dropdown_value():
    if model_manager.get_loaded_model_config():
        return model_manager.get_loaded_model_config().get("display_name")
    else:
        return model_manager.get_model_list()[0] if len(model_manager.get_model_list()) > 0 else None


def update_model_selector_choices():
    return gr.update(choices=model_manager.get_model_list(), value=update_select_model_dropdown_value())


def add_model(model_name: Optional[str], original_repo: str, mlx_repo: str, quantize: str, default_language: str, default_system_prompt: Optional[str], multimodal_ability: List[str]):
    try:
        model_manager.add_config(original_repo, mlx_repo, model_name, quantize, default_language, default_system_prompt, multimodal_ability)
    except Exception as e:
        raise gr.Error(str(e))


def add_api_model(model_name: str, api_key: str, nick_name: Optional[str] = None, base_url: Optional[str] = None, system_prompt: Optional[str] = None):
    try:
        model_manager.add_api_config(model_name, api_key, nick_name=nick_name, base_url=base_url, system_prompt=system_prompt)
    except Exception as e:
        raise gr.Error(str(e))


def update_slider_config(slider_new_min: Union[int, float], slider_new_max: Union[int, float], slider_value: Union[int, float]):
    if not slider_new_min or not slider_new_max:
        return gr.update()

    if slider_new_min > slider_new_max:
        return gr.update()

    value_to_set = slider_value if slider_value is not None else (slider_new_min + slider_new_max) / 2

    if isinstance(slider_new_min, int) and isinstance(slider_new_max, int):
        value_to_set = int(value_to_set)

    if value_to_set < slider_new_min:
        value_to_set = slider_new_min

    if value_to_set > slider_new_max:
        value_to_set = slider_new_max

    return gr.update(minimum=slider_new_min, maximum=slider_new_max, value=value_to_set)


def update_model_max_length(slider_value: Union[int, float]):
    try:
        model = get_loaded_model()
        if model is not None:
            if isinstance(model, OpenAIModel):
                return update_slider_config(1, 1048576, slider_value)
            try:
                model_max_length = model.tokenizer.model_max_length
                if model_max_length is None or model_max_length <= 0:
                    model_max_length = 32768
            except Exception as e:
                model_max_length = 32768
                logger.error("Error while updating model max length.")
            return update_slider_config(1, model_max_length, slider_value)
        else:
            return gr.update()
    except RuntimeError:
        return gr.update()
    except Exception as e:
        raise gr.Error(str(e))


def bytes_to_gigabytes(value):
    return value / 1024 ** 3


def get_memory_usage() -> str:
    memory_usage_bytes = model_manager.get_system_memory_usage()
    total_memory_bytes = model_manager.get_device_info()["memory_size"]

    memory_usage_gb = bytes_to_gigabytes(memory_usage_bytes) if isinstance(memory_usage_bytes, (int, float)) else "N/A"
    total_memory_gb = bytes_to_gigabytes(total_memory_bytes) if isinstance(total_memory_bytes, (int, float)) else "N/A"

    return "{:.2f} GB | {:.2f} GB".format(memory_usage_gb, total_memory_gb)


def update_all_memory_usage() -> list[str]:
    memory_usage = get_memory_usage()
    return [memory_usage, memory_usage]


def get_rag_enabled_status() -> bool:
    return rag_manager.is_enabled()


def get_rag_status() -> str:
    return rag_manager.get_status()[1]


def upload_and_index_file(files) -> Tuple[str, str]:
    if not files:
        return "No files selected.", get_rag_status()

    results = []
    for file_path in files:
        try:
            file_content = file_manager.load_file(Path(file_path), raw_content_only=True)
            if file_content:
                success, message = rag_manager.add_document(Path(file_path), file_content)
                results.append(f"{Path(file_path).name}: {message}")
            else:
                results.append(f"{Path(file_path).name}: Failed to load content")
        except Exception as e:
            results.append(f"{Path(file_path).name}: Error - {str(e)}")

    return "\n".join(results), get_rag_status()


def clear_rag_index() -> Tuple[str, str]:
    try:
        rag_manager.clear_index()
        return "RAG index cleared successfully.", get_rag_status()
    except Exception as e:
        return f"Error clearing RAG index: {str(e)}", get_rag_status()


def toggle_rag_enabled(enabled: bool) -> str:
    if enabled:
        rag_manager.enable()
        return "RAG enabled."
    else:
        rag_manager.disable()
        return "RAG disabled."


def update_rag_parameters(chunk_size: int, chunk_overlap: int, similarity_threshold: float) -> str:
    try:
        updated = rag_manager.update_parameters(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            similarity_threshold=similarity_threshold
        )

        if updated:
            return "RAG parameters updated. Consider re-indexing documents for optimal performance."
        else:
            return "RAG parameters updated."
    except Exception as e:
        return f"Error updating RAG parameters: {str(e)}"


def get_rag_parameters() -> Tuple[int, int, int, float]:
    params = rag_manager.get_parameters()
    return (
        params["chunk_size"],
        params["chunk_overlap"],
        params["n_results"],
        params["similarity_threshold"]
    )


def enhance_message_with_rag(message_text: str, rag_enabled: bool, n_results: int = 5) -> str:
    if not rag_enabled:
        return message_text

    try:
        retrieved_docs = rag_manager.retrieve(message_text, n_results=n_results)

        if retrieved_docs:
            enhanced_message = """
            Based on the following relevant documents: {}

            ---

            User question: {}
            """.format(retrieved_docs, message_text)
            return enhanced_message
        else:
            return message_text
    except Exception as e:
        logger.error(f"Error in RAG retrieval: {e}")
        return message_text


def create_slider(min_val, max_val, default_val, label_key, **kwargs):
    return gr.Slider(
        minimum=min_val,
        maximum=max_val,
        value=default_val,
        label=get_text(label_key),
        render=False,
        interactive=True,
        **kwargs
    )


def create_textbox(label_key, placeholder_key=None, **kwargs):
    params = {
        'label': get_text(label_key),
        'render': False,
        'interactive': True,
        **kwargs
    }
    if placeholder_key:
        params['placeholder'] = get_text(placeholder_key)
    return gr.Textbox(**params)


def create_model_controls():
    memory_usage = gr.Textbox(
        label=get_text("Page.Chat.SystemStatusBlock.Textbox.memory_usage.label"),
        interactive=False,
        render=False
    )

    model_selector = gr.Dropdown(
        label=get_text("Page.Chat.LoadModelBlock.Dropdown.model_selector.label"),
        choices=model_manager.get_model_list(),
        render=False,
        interactive=True
    )

    model_status = gr.Textbox(
        value=get_load_model_status,
        show_label=False,
        render=False,
        interactive=False
    )

    load_button = gr.Button(
        value=get_text("Page.Chat.LoadModelBlock.Button.load_model.value"),
        render=False,
        interactive=True
    )

    return memory_usage, model_selector, model_status, load_button


def create_generation_params():
    slider_configs = {
        'temperature': (0.0, 2.0, 0.6),
        'top_k': (0, 100, 20),
        'top_p': (0.0, 1.0, 0.95),
        'min_p': (0.0, 1.0, 0.0),
        'max_tokens': (1, 32768, 4096),
        'repetition_penalty': (0.0, 2.0, 1.0)
    }

    sliders = {}
    for param, (min_val, max_val, default) in slider_configs.items():
        sliders[param] = create_slider(
            min_val, max_val, default,
            f"Page.Chat.Accordion.AdvancedSetting.Slider.{param}.label"
        )

    return sliders


def create_rag_params():
    chunk_size, chunk_overlap, n_results, similarity_threshold = get_rag_parameters()

    rag_params = {
        'chunk_size': gr.Slider(
            minimum=100,
            maximum=2000,
            value=chunk_size,
            step=50,
            label=get_text("Page.Chat.Accordion.RAGSetting.Slider.chunk_size.label"),
            render=False,
            interactive=True
        ),
        'chunk_overlap': gr.Slider(
            minimum=0,
            maximum=500,
            value=chunk_overlap,
            step=10,
            label=get_text("Page.Chat.Accordion.RAGSetting.Slider.chunk_overlap.label"),
            render=False,
            interactive=True
        ),
        'n_results': gr.Slider(
            minimum=1,
            maximum=10,
            value=n_results,
            step=1,
            label=get_text("Page.Chat.Accordion.RAGSetting.Slider.n_results.label"),
            render=False,
            interactive=True
        ),
        'similarity_threshold': gr.Slider(
            minimum=0.0,
            maximum=1.0,
            value=similarity_threshold,
            step=0.05,
            label=get_text("Page.Chat.Accordion.RAGSetting.Slider.similarity_threshold.label"),
            render=False,
            interactive=True
        )
    }

    return rag_params


def setup_model_sync_events(chat_selector, completion_selector, chat_load_btn, completion_load_btn,
                            chat_status, completion_status, chat_system_prompt, chat_max_tokens, completion_max_tokens):
    chat_selector.select(
        fn=lambda x: x,
        inputs=[chat_selector],
        outputs=[completion_selector]
    )

    completion_selector.select(
        fn=lambda x: x,
        inputs=[completion_selector],
        outputs=[chat_selector]
    )

    chat_load_btn.click(
        fn=chat_load_model_callback,
        inputs=[chat_selector],
        outputs=[chat_status, chat_system_prompt]
    ).then(
        fn=lambda x: x,
        inputs=[chat_status],
        outputs=[completion_status]
    ).then(
        fn=update_model_max_length,
        inputs=[chat_max_tokens],
        outputs=[chat_max_tokens]
    ).then(
        fn=update_model_max_length,
        inputs=[completion_max_tokens],
        outputs=[completion_max_tokens]
    )

    completion_load_btn.click(
        fn=completion_load_model_callback,
        inputs=[completion_selector],
        outputs=[completion_status]
    ).then(
        fn=lambda x: x,
        inputs=[completion_status],
        outputs=[chat_status]
    ).then(
        fn=update_model_max_length,
        inputs=[chat_max_tokens],
        outputs=[chat_max_tokens]
    ).then(
        fn=update_model_max_length,
        inputs=[completion_max_tokens],
        outputs=[completion_max_tokens]
    )


def setup_model_management_events(local_form, api_form, model_list, chat_selector, completion_selector):
    local_form['add_button'].click(
        fn=add_model,
        inputs=[
            local_form['model_name'],
            local_form['original_repo'],
            local_form['mlx_repo'],
            local_form['quantize'],
            local_form['default_language'],
            local_form['system_prompt'],
            local_form['multimodal']
        ]
    ).then(
        fn=update_model_management_models_list,
        outputs=[model_list]
    ).then(
        fn=update_model_selector_choices,
        outputs=[chat_selector]
    ).then(
        fn=update_model_selector_choices,
        outputs=[completion_selector]
    )

    api_form['add_button'].click(
        fn=add_api_model,
        inputs=[
            api_form['model_name'],
            api_form['api_key'],
            api_form['nick_name'],
            api_form['base_url'],
            api_form['system_prompt']
        ]
    ).then(
        fn=update_model_management_models_list,
        outputs=[model_list]
    ).then(
        fn=update_model_selector_choices,
        outputs=[chat_selector]
    ).then(
        fn=update_model_selector_choices,
        outputs=[completion_selector]
    )


def clear_cache():
    stop_generation()
    file_manager.clear()
    model_manager.close_active_generator()


with gr.Blocks(fill_height=True, fill_width=True, title="Chat with MLX") as app:
    update_memory_usage_timer = gr.Timer(value=1, active=True)

    chat_memory, chat_model_selector, chat_model_status, chat_load_button = create_model_controls()
    chat_params = create_generation_params()
    rag_params = create_rag_params()

    chat_system_prompt_textbox = gr.Textbox(
        label=get_text("Page.Chat.ChatSystemPromptBlock.Textbox.system_prompt.label"),
        placeholder=get_text("Page.Chat.ChatSystemPromptBlock.Textbox.system_prompt.placeholder"),
        value=model_manager.get_system_prompt,
        lines=3,
        max_lines=5,
        show_copy_button=True,
        render=False,
        scale=9
    )

    chat_default_system_prompt_button = gr.Button(
        value=get_text("Page.Chat.ChatSystemPromptBlock.Button.default_system_prompt.value"),
        render=False,
        scale=1
    )

    chat_rag_form = {
        'rag_enabled': gr.Checkbox(
            label=get_text("Page.Chat.Accordion.RAGSetting.Checkbox.rag_enabled.label"),
            value=get_rag_enabled_status,
            interactive=True,
            render=False
        ),
        'file_upload': gr.File(
            label=get_text("Page.Chat.Accordion.RAGSetting.File.file_upload.label"),
            file_count="multiple",
            file_types=[".txt", ".pdf", ".docx", ".pptx", ".xlsx", ".xls", ".md", ".csv"],
            render=False
        ),
        'upload_button': gr.Button(
            value=get_text("Page.Chat.Accordion.RAGSetting.Button.upload.value"),
            render=False
        ),
        'clear_button': gr.Button(
            value=get_text("Page.Chat.Accordion.RAGSetting.Button.clear.value"),
            render=False
        ),
        'upload_status': create_textbox("Page.Chat.Accordion.RAGSetting.Textbox.upload_status.label", interactive=False, render=False),
        'rag_status': create_textbox("Page.Chat.Accordion.RAGSetting.Textbox.rag_status.label", interactive=False, render=False),
        'update_params_button': gr.Button(
            value=get_text("Page.Chat.Accordion.RAGSetting.Button.update_params.value"),
            render=False
        ),
        'params_status': create_textbox("Page.Chat.Accordion.RAGSetting.Textbox.params_status.label", interactive=False, render=False)
    }

    completion_memory, completion_model_selector, completion_model_status, completion_load_button = create_model_controls()
    completion_params = create_generation_params()

    local_model_form = {
        'model_name': create_textbox("Page.ModelManagement.AddLocalModelBlock.Textbox.model_name.label",
                                     "Page.ModelManagement.AddLocalModelBlock.Textbox.model_name.placeholder"),
        'original_repo': create_textbox("Page.ModelManagement.AddLocalModelBlock.Textbox.original_repo.label",
                                        "Page.ModelManagement.AddLocalModelBlock.Textbox.original_repo.placeholder"),
        'mlx_repo': create_textbox("Page.ModelManagement.AddLocalModelBlock.Textbox.mlx_repo.label",
                                   "Page.ModelManagement.AddLocalModelBlock.Textbox.mlx_repo.placeholder"),
        'quantize': gr.Dropdown(
            label=get_text("Page.ModelManagement.AddLocalModelBlock.Dropdown.quantize.label"),
            choices=["None", "4bit", "8bit", "bf16", "bf32"],
            value="None",
            interactive=True,
            render=False
        ),
        'default_language': gr.Dropdown(
            label=get_text("Page.ModelManagement.AddLocalModelBlock.Dropdown.default_language.label"),
            choices=["multi"],
            interactive=True,
            render=False
        ),
        'system_prompt': create_textbox("Page.ModelManagement.AddLocalModelBlock.Textbox.default_system_prompt.label"),
        'multimodal': gr.Dropdown(
            label=get_text("Page.ModelManagement.AddLocalModelBlock.Dropdown.multimodal_ability.label"),
            choices=["None", "vision"],
            value="None",
            multiselect=True,
            render=False
        ),
        'add_button': gr.Button(
            value=get_text("Page.ModelManagement.AddLocalModelBlock.Button.add.value"),
            render=False
        )
    }

    api_model_form = {
        'model_name': create_textbox("Page.ModelManagement.AddAPIModelBlock.Textbox.model_name.label",
                                     "Page.ModelManagement.AddAPIModelBlock.Textbox.model_name.placeholder"),
        'nick_name': create_textbox("Page.ModelManagement.AddAPIModelBlock.Textbox.nick_name.label"),
        'api_key': create_textbox("Page.ModelManagement.AddAPIModelBlock.Textbox.api_key.label",
                                  "Page.ModelManagement.AddAPIModelBlock.Textbox.api_key.placeholder"),
        'base_url': create_textbox("Page.ModelManagement.AddAPIModelBlock.Textbox.base_url.label",
                                   "Page.ModelManagement.AddAPIModelBlock.Textbox.base_url.placeholder"),
        'system_prompt': create_textbox("Page.ModelManagement.AddLocalModelBlock.Textbox.default_system_prompt.label"),
        'add_button': gr.Button(
            value=get_text("Page.ModelManagement.AddLocalModelBlock.Button.add.value"),
            render=False
        )
    }

    model_list = gr.Dataframe(
        headers=[get_text("Page.ModelManagement.Dataframe.model_list.headers")],
        value=update_model_management_models_list(),
        datatype="str",
        row_count=(10, "dynamic"),
        render=False,
        interactive=False
    )

    update_memory_usage_timer.tick(
        fn=update_all_memory_usage,
        outputs=[chat_memory, completion_memory]
    )

    gr.HTML("<h2>Chat with MLX</h2>")

    with gr.Tab(get_text("Tab.chat")):
        with gr.Row():
            with gr.Column(scale=2):
                with gr.Row():
                    chat_memory.render()

                with gr.Row():
                    gr.Markdown(f"## {get_text('Page.Chat.Markdown.configuration')}")

                    chat_model_selector.render()
                    chat_model_status.render()
                    chat_load_button.render()

                with gr.Accordion(label=get_text("Page.Chat.Accordion.AdvancedSetting.label"), open=False):
                    with gr.Group():
                        for slider in chat_params.values():
                            slider.render()

                with gr.Accordion(label=get_text("Page.Chat.Accordion.RAGSetting.label"), open=False):
                    chat_rag_form['rag_enabled'].render()
                    chat_rag_form['rag_status'].render()

                    with gr.Row():
                        chat_rag_form['file_upload'].render()

                    with gr.Row():
                        chat_rag_form['upload_button'].render()
                        chat_rag_form['clear_button'].render()

                    chat_rag_form['upload_status'].render()

                    with gr.Group():
                        with gr.Row():
                            rag_params['chunk_size'].render()
                            rag_params['chunk_overlap'].render()

                        with gr.Row():
                            rag_params['n_results'].render()
                            rag_params['similarity_threshold'].render()

                    chat_rag_form['update_params_button'].render()
                    chat_rag_form['params_status'].render()

            with gr.Column(scale=8):
                with gr.Row(equal_height=True):
                    chat_system_prompt_textbox.render()
                    chat_default_system_prompt_button.render()

                    chat_default_system_prompt_button.click(
                        fn=get_default_system_prompt_callback,
                        outputs=[chat_system_prompt_textbox]
                    )
                    chat_system_prompt_textbox.change(
                        fn=model_manager.set_custom_prompt,
                        inputs=[chat_system_prompt_textbox]
                    )

                chatbot = gr.Chatbot(
                    type="messages",
                    show_copy_button=True,
                    render=False,
                    latex_delimiters=[
                        {"left": "\\begin{equation}", "right": "\\end{equation}", "display": True},
                        {"left": "\\begin{equation*}", "right": "\\end{equation*}", "display": True},
                        {"left": "\\begin{align}", "right": "\\end{align}", "display": True},
                        {"left": "\\begin{align*}", "right": "\\end{align*}", "display": True},
                        {"left": "\\begin{alignat}", "right": "\\end{alignat}", "display": True},
                        {"left": "\\begin{alignat*}", "right": "\\end{alignat*}", "display": True},
                        {"left": "\\begin{gather}", "right": "\\end{gather}", "display": True},
                        {"left": "\\begin{gather*}", "right": "\\end{gather*}", "display": True},
                        {"left": "\\begin{eqnarray}", "right": "\\end{eqnarray}", "display": True},
                        {"left": "\\begin{eqnarray*}", "right": "\\end{eqnarray*}", "display": True},
                        {"left": "\\begin{multline}", "right": "\\end{multline}", "display": True},
                        {"left": "\\begin{multline*}", "right": "\\end{multline*}", "display": True},
                        {"left": "\\begin{split}", "right": "\\end{split}", "display": True},
                        {"left": "\\begin{cases}", "right": "\\end{cases}", "display": True},
                        {"left": "\\begin{matrix}", "right": "\\end{matrix}", "display": True},
                        {"left": "\\begin{pmatrix}", "right": "\\end{pmatrix}", "display": True},
                        {"left": "\\begin{bmatrix}", "right": "\\end{bmatrix}", "display": True},
                        {"left": "\\begin{vmatrix}", "right": "\\end{vmatrix}", "display": True},
                        {"left": "\\begin{Vmatrix}", "right": "\\end{Vmatrix}", "display": True},
                        {"left": "\\begin{CD}", "right": "\\end{CD}", "display": True},
                        {"left": "\\[", "right": "\\]", "display": True},
                        {"left": "$$", "right": "$$", "display": True},
                        {"left": "\\(", "right": "\\)", "display": False},
                        {"left": "$", "right": "$", "display": False}
                    ]
                )

                chatbot.clear(
                    fn=clear_cache
                )

                gr.ChatInterface(
                    multimodal=True,
                    chatbot=chatbot,
                    type="messages",
                    fn=managed_chat_generator,
                    title=None,
                    autofocus=False,
                    fill_height=True,
                    fill_width=True,
                    save_history=True,
                    additional_inputs=[chat_system_prompt_textbox] + list(chat_params.values()) + [chat_rag_form['rag_enabled'], rag_params['n_results']]
                )

    with gr.Tab(get_text("Tab.completion"), interactive=True):
        with gr.Row():
            with gr.Column(scale=2):
                completion_memory.render()

                gr.Markdown(f"## {get_text('Page.Chat.Markdown.configuration')}")

                completion_model_selector.render()
                completion_model_status.render()
                completion_load_button.render()

                with gr.Row(visible=False):
                    with gr.Accordion(label=get_text("Page.Chat.Accordion.AdvancedSetting.label"), open=True):
                        for slider in completion_params.values():
                            slider.render()

            with gr.Column(scale=8):
                completion_interface = gr.Interface(
                    clear_btn=None,
                    flagging_mode="never",
                    fn=managed_completion_generator,
                    inputs=[
                               gr.Textbox(lines=10, show_copy_button=True, render=True,
                                          label=get_text("Page.Completion.Textbox.prompt.label")),
                           ] + list(completion_params.values()),
                    outputs=[
                        gr.Textbox(lines=25, show_copy_button=True, render=True,
                                   label=get_text("Page.Completion.Textbox.output.label"))
                    ],
                    submit_btn=get_text("Page.Completion.Button.submit.value"),
                    stop_btn=get_text("Page.Completion.Button.stop.value"),
                )

    with gr.Tab(get_text("Tab.model_management"), interactive=True):
        with gr.Row(equal_height=True):
            with gr.Column(scale=5):
                model_list.render()

            with gr.Column(scale=5):
                with gr.Tab(get_text("Page.ModelManagement.Tab.local_model")):
                    for component in local_model_form.values():
                        component.render()

                with gr.Tab(get_text("Page.ModelManagement.Tab.openai_api")):
                    for component in api_model_form.values():
                        component.render()

    setup_model_sync_events(
        chat_model_selector, completion_model_selector,
        chat_load_button, completion_load_button,
        chat_model_status, completion_model_status,
        chat_system_prompt_textbox,
        chat_params['max_tokens'], completion_params['max_tokens']
    )

    setup_model_management_events(
        local_model_form, api_model_form, model_list,
        chat_model_selector, completion_model_selector
    )

    chat_rag_form['rag_enabled'].change(
        fn=toggle_rag_enabled,
        inputs=[chat_rag_form['rag_enabled']],
        outputs=[chat_rag_form['upload_status']]
    )

    chat_rag_form['upload_button'].click(
        fn=upload_and_index_file,
        inputs=[chat_rag_form['file_upload']],
        outputs=[chat_rag_form['upload_status'], chat_rag_form['rag_status']]
    )

    chat_rag_form['clear_button'].click(
        fn=clear_rag_index,
        outputs=[chat_rag_form['upload_status'], chat_rag_form['rag_status']]
    )

    chat_rag_form['update_params_button'].click(
        fn=update_rag_parameters,
        inputs=[
            rag_params['chunk_size'],
            rag_params['chunk_overlap'],
            rag_params['similarity_threshold']
        ],
        outputs=[chat_rag_form['params_status']]
    )

    app.load(
        fn=update_model_management_models_list,
        outputs=[model_list]
    ).then(
        fn=update_model_selector_choices,
        outputs=[chat_model_selector]
    ).then(
        fn=update_model_selector_choices,
        outputs=[completion_model_selector]
    ).then(
        fn=update_model_max_length,
        inputs=[chat_params['max_tokens']],
        outputs=[chat_params['max_tokens']]
    ).then(
        fn=update_model_max_length,
        inputs=[completion_params['max_tokens']],
        outputs=[completion_params['max_tokens']]
    ).then(
        fn=update_all_memory_usage,
        outputs=[chat_memory, completion_memory]
    ).then(
        fn=get_rag_status,
        outputs=[chat_rag_form['rag_status']]
    )


def exit_handler():
    model_manager.close_model()


atexit.register(exit_handler)


def start(port: int, share: bool = False, in_browser: bool = True) -> None:
    logger.info(f"Starting the app on port {port} with share={share} and in_browser={in_browser}")
    app.launch(server_port=port, inbrowser=in_browser, share=share, pwa=True)


def main():
    parser = argparse.ArgumentParser(description="Chat with MLX")
    parser.add_argument(
        "--port",
        type=int,
        default=7860,
        help="The port number to run the application on (default: 7860)"
    )
    parser.add_argument(
        "--share",
        action="store_true",
        help="Enable sharing the application link externally"
    )
    parser.add_argument(
        "--no-browser",
        action="store_true",
        help="Do not open the application in the default web browser"
    )
    args = parser.parse_args()

    start(port=args.port, share=args.share, in_browser=not args.no_browser)


if __name__ == "__main__":
    main()
